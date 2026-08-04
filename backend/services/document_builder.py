import json
import os
import re
import subprocess
import sys
import tempfile
from typing import Any
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet

# ── Path to the Node.js PPT generator script ─────────────────────────────────
# Place generate_pptx.js in the same directory as this file, or set an env var.
_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PPTX_SCRIPT = os.environ.get(
    "PPTX_SCRIPT_PATH",
    os.path.join(_SCRIPT_DIR, "generate_pptx.js"),
)


# ── DOCUMENT TYPE DETECTOR ───────────────────────────────────────────────────
def detect_document_type(context: str, question: str) -> str:
    """
    Classify the source document so the AI and PPT builder can choose
    the right visual language.
    """
    combined = (context + " " + question).lower()

    if any(k in combined for k in ("order", "tariff order", "merc", "cerc", "regulatory")):
        return "regulatory_order"
    if any(k in combined for k in ("roadmap", "mission", "hydrogen", "trajectory", "strategy")):
        return "roadmap"
    if any(k in combined for k in ("scheme", "incentive", "subsidy", "grant", "benefit")):
        return "incentive_scheme"
    if any(k in combined for k in ("policy", "act", "regulation", "gazette", "notification")):
        return "policy"
    return "generic"


# ── JSON UTILITIES ───────────────────────────────────────────────────────────
def parse_json_response(raw_text: str) -> dict[str, Any]:
    raw = raw_text.strip()
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"```$", "", raw.strip())

    start = raw.find("{")
    end = raw.rfind("}")
    if start == -1 or end == -1:
        raise ValueError("Could not locate a JSON object in Claude response.")

    return json.loads(raw[start:end + 1])


# ── DOCX ─────────────────────────────────────────────────────────────────────
def _truncate_text(text: str, max_chars: int) -> str:
    if not text:
        return ""
    stripped = text.strip()
    if len(stripped) <= max_chars:
        return stripped
    return stripped[:max_chars].rsplit(" ", 1)[0] + "..."


def _chunk_list(items: list, chunk_size: int) -> list[list]:
    return [items[i:i + chunk_size] for i in range(0, len(items), chunk_size)]


def normalize_document_content(content: dict[str, Any]) -> dict[str, Any]:
    if not isinstance(content, dict):
        raise ValueError("Document content must be a JSON object.")

    title = str(content.get("title", "Policy Document")).strip() or "Policy Document"
    subtitle = str(content.get("subtitle", "")).strip()
    sections = content.get("sections")

    if not isinstance(sections, list) or not sections:
        sections = [{
            "heading": "Key Points",
            "summary": str(content.get("summary", "")).strip() or "Summary of the document.",
            "bullets": [str(b).strip() for b in content.get("bullets", []) if str(b).strip()],
        }]

    normalized_sections = []
    for section in sections:
        if not isinstance(section, dict):
            continue
        normalized_sections.append({
            "heading": str(section.get("heading", "Section")).strip() or "Section",
            "summary": str(section.get("summary", "")).strip(),
            "bullets": [str(b).strip() for b in section.get("bullets", []) if str(b).strip()],
        })

    if not normalized_sections:
        normalized_sections = [{
            "heading": "Key Points",
            "summary": str(content.get("summary", "")).strip() or "Summary of the document.",
            "bullets": [str(b).strip() for b in content.get("bullets", []) if str(b).strip()],
        }]

    return {
        "title": title,
        "subtitle": subtitle,
        "sections": normalized_sections,
        "key_stats": [str(s).strip() for s in content.get("key_stats", []) if str(s).strip()],
        "conclusion": str(content.get("conclusion", "")).strip(),
    }


def build_docx(content: dict[str, Any]) -> str:
    from docx import Document

    data = normalize_document_content(content)
    doc = Document()
    doc.add_heading(data["title"], level=0)
    if data["subtitle"]:
        doc.add_paragraph(data["subtitle"])
    for section in data["sections"]:
        doc.add_heading(section["heading"], level=1)
        if section["summary"]:
            doc.add_paragraph(section["summary"])
        for bullet in section["bullets"]:
            doc.add_paragraph(bullet, style="List Bullet")
    if data["key_stats"]:
        doc.add_heading("Key Statistics", level=1)
        for stat in data["key_stats"]:
            doc.add_paragraph(stat, style="List Bullet")
    if data["conclusion"]:
        doc.add_heading("Conclusion", level=1)
        doc.add_paragraph(data["conclusion"])

    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    temp_file.close()
    doc.save(temp_file.name)
    return temp_file.name


# ── PDF ───────────────────────────────────────────────────────────────────────
def build_pdf(content: dict[str, Any]) -> str:
    data = normalize_document_content(content)
    path = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    path.close()

    doc = SimpleDocTemplate(path.name, pagesize=A4)
    styles = getSampleStyleSheet()
    story = [Paragraph(data["title"], styles["Title"]), Spacer(1, 12)]
    if data["subtitle"]:
        story += [Paragraph(data["subtitle"], styles["Heading2"]), Spacer(1, 12)]
    for section in data["sections"]:
        story.append(Paragraph(section["heading"], styles["Heading2"]))
        story.append(Paragraph(section["summary"], styles["Normal"]))
        story.append(Spacer(1, 8))
        for bullet in section["bullets"]:
            story.append(Paragraph(f"• {bullet}", styles["Normal"]))
        story.append(Spacer(1, 12))
    if data["key_stats"]:
        story.append(Paragraph("Key Statistics", styles["Heading2"]))
        for stat in data["key_stats"]:
            story.append(Paragraph(f"• {stat}", styles["Normal"]))
        story.append(Spacer(1, 12))
    if data["conclusion"]:
        story.append(Paragraph("Conclusion", styles["Heading2"]))
        story.append(Paragraph(data["conclusion"], styles["Normal"]))

    doc.build(story)
    return path.name


# ── PPTX — DYNAMIC AI-DRIVEN BUILDER ─────────────────────────────────────────
# `content` is now {title, narrative_label, theme, slides[]} where Claude has
# already chosen fonts/colors/layout sequence per document (see export.py's
# _build_ppt_prompt). No more document_type -> fixed palette/template routing
# here — the theme travels inside the JSON itself.

def build_pptx(content: dict[str, Any]) -> str:
    """
    Build a premium presentation using the pptxgenjs Node.js script.
    Falls back to a plain-text python-pptx version only if Node.js isn't
    available in the deployment environment (no charts/theme in that path —
    keep Node available in production so the dynamic renderer is always used).
    """
    try:
        return _build_pptx_node(content)
    except Exception as e:
        print(f"[PPTX] Node builder failed ({e}), falling back to python-pptx", file=sys.stderr)
        return _build_pptx_python_fallback(content)


def _build_pptx_node(content: dict[str, Any]) -> str:
    """Write JSON to temp file, call Node.js generator, return output path."""
    if not os.path.exists(PPTX_SCRIPT):
        raise FileNotFoundError(f"PPTX script not found: {PPTX_SCRIPT}")

    node_exec = _find_node()
    if not node_exec:
        raise RuntimeError("Node.js not found in PATH")

    input_file = tempfile.NamedTemporaryFile(delete=False, suffix=".json", mode="w")
    json.dump(content, input_file, ensure_ascii=False)
    input_file.close()

    output_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    output_file.close()

    try:
        result = subprocess.run(
            [node_exec, PPTX_SCRIPT, input_file.name, output_file.name],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode != 0:
            raise RuntimeError(f"Node script error: {result.stderr[:500]}")
        stdout = result.stdout.strip()
        if not stdout.startswith("DONE:"):
            raise RuntimeError(f"Unexpected node output: {stdout[:200]}")
        return output_file.name
    finally:
        os.unlink(input_file.name)


def _find_node() -> str | None:
    import shutil
    for candidate in ["node", "nodejs"]:
        path = shutil.which(candidate)
        if path:
            return path
    return None


def _build_pptx_python_fallback(content: dict[str, Any]) -> str:
    """
    Minimal plain-text fallback — only used if Node.js is unavailable.
    No theme colors/fonts and no charts here (that logic lives in
    generate_pptx.js); this just makes sure an export never hard-fails.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(blank)
    tb = title_slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    tb.text_frame.text = content.get("title", "Policy Presentation")
    tb.text_frame.paragraphs[0].font.size = None  # keep default; theme not available in fallback

    for sd in content.get("slides", []):
        slide = prs.slides.add_slide(blank)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_box.text_frame.text = sd.get("title") or sd.get("eyebrow") or ""

        body = list(sd.get("bullets") or sd.get("takeaways") or [])
        for it in sd.get("items", []) or []:
            if isinstance(it, dict):
                body.append(f"{it.get('title', '')}: {it.get('detail', '') or it.get('text', '')}")
        for ev in sd.get("events", []) or []:
            if isinstance(ev, dict):
                body.append(f"{ev.get('date', ev.get('year', ''))}: {ev.get('label', '')}")

        if body:
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))
            tf = box.text_frame
            tf.word_wrap = True
            for i, b in enumerate(body[:10]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {_truncate_text(str(b), 150)}"

    return _save_pptx(prs)


def _save_pptx(prs) -> str:
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    temp_file.close()
    prs.save(temp_file.name)
    return temp_file.name